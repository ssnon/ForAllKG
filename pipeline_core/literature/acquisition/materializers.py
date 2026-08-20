from __future__ import annotations

import csv
import hashlib
import html
import json
import os
import re
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Protocol

from pipeline_core.literature.acquisition.materialization_contracts import MaterializationPolicy


@dataclass(frozen=True)
class MaterializerOutput:
    materializer: str
    markdown: str
    asset_source_dir: Path | None = None
    primary_markdown_source: Path | None = None


class DocumentMaterializer(Protocol):
    materializer_id: str

    def materialize(
        self,
        *,
        source_path: Path,
        policy: MaterializationPolicy,
    ) -> MaterializerOutput:
        ...


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _markdown_escape(value: object) -> str:
    text = "" if value is None else str(value)
    text = text.replace("\\", "\\\\")
    text = text.replace("|", "\\|")
    text = text.replace("\r", " ").replace("\n", " ")
    return text.strip()


class MarkerPdfMaterializer:
    materializer_id = "marker_cli_pdf_v1"

    def materialize(
        self,
        *,
        source_path: Path,
        policy: MaterializationPolicy,
    ) -> MaterializerOutput:
        if not policy.materialize_pdf:
            raise RuntimeError("pdf_materialization_disabled")
        command_path = shutil.which(policy.marker_command)
        if command_path is None:
            raise RuntimeError(
                f"marker_command_not_found:{policy.marker_command}"
            )

        with tempfile.TemporaryDirectory(prefix="graphagents_marker_") as tmp:
            output_root = Path(tmp) / "output"
            output_root.mkdir(parents=True, exist_ok=True)
            command = [
                command_path,
                str(source_path),
                "--output_dir",
                str(output_root),
                *policy.marker_extra_args,
            ]
            marker_env = os.environ.copy()
            for key in policy.marker_environment_unset:
                marker_env.pop(str(key), None)
            for key, value in policy.marker_environment_overrides.items():
                marker_env[str(key)] = str(value)

            completed = subprocess.run(
                command,
                text=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                timeout=policy.marker_timeout_seconds,
                env=marker_env,
            )
            if completed.returncode != 0:
                tail = completed.stdout[-4000:]
                raise RuntimeError(
                    "marker_failed:"
                    f"exit={completed.returncode}:"
                    f"{tail}"
                )

            candidates = [
                path
                for path in output_root.rglob("*.md")
                if path.is_file()
            ]
            if not candidates:
                raise RuntimeError("marker_produced_no_markdown")

            # Marker normally produces one Markdown result. If ancillary
            # Markdown exists, choose the largest non-empty candidate and keep
            # its parent directory as the asset-relative package root.
            candidates = sorted(
                candidates,
                key=lambda path: (
                    -path.stat().st_size,
                    str(path),
                ),
            )
            selected = candidates[0]
            markdown = selected.read_text(
                encoding="utf-8",
                errors="replace",
            ).strip()
            if not markdown:
                raise RuntimeError("marker_produced_empty_markdown")

            # Persist the temporary Marker parent to another temp-like stable
            # directory owned by the caller through a copied tree snapshot.
            snapshot = Path(
                tempfile.mkdtemp(prefix="graphagents_marker_snapshot_")
            )
            for child in selected.parent.iterdir():
                target = snapshot / child.name
                if child.is_dir():
                    shutil.copytree(child, target)
                else:
                    shutil.copy2(child, target)

            primary = snapshot / selected.name
            return MaterializerOutput(
                materializer=self.materializer_id,
                markdown=markdown,
                asset_source_dir=snapshot,
                primary_markdown_source=primary,
            )


class TextMaterializer:
    materializer_id = "plain_text_v1"

    def materialize(
        self,
        *,
        source_path: Path,
        policy: MaterializationPolicy,
    ) -> MaterializerOutput:
        if not policy.materialize_txt:
            raise RuntimeError("text_materialization_disabled")
        if source_path.stat().st_size > policy.text_max_bytes:
            raise RuntimeError("text_source_exceeds_limit")
        text = source_path.read_text(
            encoding="utf-8",
            errors="replace",
        ).strip()
        if not text:
            raise RuntimeError("empty_text_source")
        return MaterializerOutput(
            materializer=self.materializer_id,
            markdown=text,
        )


class CsvMaterializer:
    materializer_id = "csv_markdown_v1"

    def materialize(
        self,
        *,
        source_path: Path,
        policy: MaterializationPolicy,
    ) -> MaterializerOutput:
        if not policy.materialize_csv:
            raise RuntimeError("csv_materialization_disabled")
        rows = []
        with source_path.open(
            "r",
            encoding="utf-8-sig",
            errors="replace",
            newline="",
        ) as handle:
            reader = csv.reader(handle)
            for index, row in enumerate(reader):
                if index >= policy.csv_max_rows:
                    break
                rows.append([_markdown_escape(value) for value in row])
        if not rows:
            raise RuntimeError("empty_csv_source")
        width = max(len(row) for row in rows)
        normalized = [
            row + [""] * (width - len(row))
            for row in rows
        ]
        header = normalized[0]
        body = normalized[1:]
        lines = [
            "# Supplementary Table",
            "",
            "| " + " | ".join(header) + " |",
            "| " + " | ".join(["---"] * width) + " |",
        ]
        lines.extend(
            "| " + " | ".join(row) + " |"
            for row in body
        )
        return MaterializerOutput(
            materializer=self.materializer_id,
            markdown="\n".join(lines).strip(),
        )


class XlsxMaterializer:
    materializer_id = "xlsx_markdown_v1"

    def materialize(
        self,
        *,
        source_path: Path,
        policy: MaterializationPolicy,
    ) -> MaterializerOutput:
        if not policy.materialize_xlsx:
            raise RuntimeError("xlsx_materialization_disabled")
        import openpyxl

        workbook = openpyxl.load_workbook(
            source_path,
            read_only=True,
            data_only=True,
        )
        sections: list[str] = []
        try:
            for sheet_index, sheet in enumerate(workbook.worksheets):
                if sheet_index >= policy.xlsx_max_sheets:
                    break
                rows = []
                for row_index, row in enumerate(
                    sheet.iter_rows(values_only=True)
                ):
                    if row_index >= policy.xlsx_max_rows_per_sheet:
                        break
                    values = list(row[: policy.xlsx_max_columns])
                    rows.append(
                        [_markdown_escape(value) for value in values]
                    )
                if not rows:
                    continue
                width = max(len(row) for row in rows)
                normalized = [
                    row + [""] * (width - len(row))
                    for row in rows
                ]
                header = normalized[0]
                body = normalized[1:]
                lines = [
                    f"# Sheet: {_markdown_escape(sheet.title)}",
                    "",
                    "| " + " | ".join(header) + " |",
                    "| " + " | ".join(["---"] * width) + " |",
                ]
                lines.extend(
                    "| " + " | ".join(row) + " |"
                    for row in body
                )
                sections.append("\n".join(lines))
        finally:
            workbook.close()
        if not sections:
            raise RuntimeError("xlsx_contains_no_materializable_sheets")
        return MaterializerOutput(
            materializer=self.materializer_id,
            markdown="\n\n".join(sections).strip(),
        )


class DocxMaterializer:
    materializer_id = "docx_mammoth_markdown_v1"

    def materialize(
        self,
        *,
        source_path: Path,
        policy: MaterializationPolicy,
    ) -> MaterializerOutput:
        if not policy.materialize_docx:
            raise RuntimeError("docx_materialization_disabled")
        import mammoth
        from markdownify import markdownify

        with source_path.open("rb") as handle:
            result = mammoth.convert_to_html(handle)
        markdown = markdownify(
            result.value,
            heading_style="ATX",
        ).strip()
        if not markdown:
            raise RuntimeError("docx_produced_empty_markdown")
        return MaterializerOutput(
            materializer=self.materializer_id,
            markdown=markdown,
        )


class PptxMaterializer:
    materializer_id = "pptx_text_markdown_v1"

    def materialize(
        self,
        *,
        source_path: Path,
        policy: MaterializationPolicy,
    ) -> MaterializerOutput:
        if not policy.materialize_pptx:
            raise RuntimeError("pptx_materialization_disabled")
        from pptx import Presentation

        presentation = Presentation(str(source_path))
        sections: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = str(shape.text or "").strip()
                    if text:
                        texts.append(text)
            if texts:
                sections.append(
                    f"# Slide {index}\n\n" + "\n\n".join(texts)
                )
        if not sections:
            raise RuntimeError("pptx_contains_no_text")
        return MaterializerOutput(
            materializer=self.materializer_id,
            markdown="\n\n".join(sections),
        )


def materializer_for(
    source_path: Path,
    policy: MaterializationPolicy,
) -> DocumentMaterializer | None:
    suffix = source_path.suffix.casefold()
    if suffix in {value.casefold() for value in policy.unsupported_extensions}:
        return None
    if suffix == ".pdf":
        return MarkerPdfMaterializer()
    if suffix == ".txt":
        return TextMaterializer()
    if suffix == ".csv":
        return CsvMaterializer()
    if suffix == ".xlsx":
        return XlsxMaterializer()
    if suffix == ".docx":
        return DocxMaterializer()
    if suffix == ".pptx":
        return PptxMaterializer()
    return None


def copy_materializer_assets(
    *,
    output: MaterializerOutput,
    package_dir: Path,
) -> int:
    if output.asset_source_dir is None:
        return 0
    count = 0
    try:
        for child in output.asset_source_dir.iterdir():
            # normalized.md is controlled by M4; keep the Marker-generated
            # original markdown only under a provenance-friendly name.
            if (
                output.primary_markdown_source is not None
                and child.resolve()
                == output.primary_markdown_source.resolve()
            ):
                target = package_dir / "marker_original.md"
            else:
                target = package_dir / child.name
            if child.is_dir():
                if target.exists():
                    shutil.rmtree(target)
                shutil.copytree(child, target)
                count += sum(1 for p in target.rglob("*") if p.is_file())
            else:
                shutil.copy2(child, target)
                count += 1
        return count
    finally:
        shutil.rmtree(output.asset_source_dir, ignore_errors=True)
